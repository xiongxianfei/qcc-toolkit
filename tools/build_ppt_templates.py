"""Build first-slice QCC PowerPoint method templates.

The Markdown files next to the generated decks remain the reviewable source notes.
This script creates static, editable PPTX decks from the catalog contract.
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from qcc_toolkit.templates import TemplateCatalogEntry, load_template_catalog

CATALOG_PATH = Path("templates/ppt/catalog.yml")
OUTPUT_ROOT = Path("templates/ppt/methods")
GENERATED_AT = datetime(2026, 7, 8, 0, 0, 0)
ZIP_TIMESTAMP = (2026, 7, 8, 0, 0, 0)


@dataclass(frozen=True)
class TemplateContent:
    """Human-facing static content for one template deck."""

    title: str
    stage_label: str
    purpose: str
    data_title: str
    data_columns: tuple[str, ...]
    data_rows: tuple[tuple[str, ...], ...]
    demo_title: str
    demo_points: tuple[str, ...]
    project_title: str
    project_points: tuple[str, ...]


@dataclass(frozen=True)
class MethodKitGuidance:
    """Required method-kit guidance for template-native decks."""

    qcc_stage_fit: tuple[str, ...]
    required_inputs: tuple[str, ...]
    when_to_use: tuple[str, ...]
    when_not_to_use: tuple[str, ...]
    interpretation_patterns: tuple[str, ...]
    common_mistakes: tuple[str, ...]
    facilitator_checklist: tuple[str, ...]
    python_assist_decision: tuple[str, ...]
    evidence_levels: tuple[str, ...]
    evidence_source_note: tuple[str, ...]


CONTENT: dict[str, TemplateContent] = {
    "pareto_chart": TemplateContent(
        title="Pareto Chart",
        stage_label="Understand Current Condition / Analyze Causes",
        purpose=(
            "Purpose: identify the vital few categories from counted project data."
        ),
        data_title=(
            "DATA ENTRY - PowerPoint Edit Data - replace categories and counts only"
        ),
        data_columns=("Defect category", "Count"),
        data_rows=(
            ("Wrong label", "42"),
            ("Missing label", "27"),
            ("Damaged barcode", "18"),
            ("Wrong carton", "9"),
            ("Other", "4"),
        ),
        demo_title="Completed demo example - DEMO EXAMPLE - not project evidence",
        demo_points=(
            "Change data only: edit the chart data table or regenerate from Python.",
            "Generated chart image slot: {{chart_image}}",
            "Calculated table reference: {{calculated_table}}",
            "Caption from evidence package: {{caption}}",
            "Data context and filters: {{data_context}}",
            "Warnings and cautions: {{warnings}}",
        ),
        project_title="Blank copyable project slide",
        project_points=(
            "Project title: [problem and period]",
            "Chart area: replace demo data or insert generated chart.",
            "Source: [system, sheet, or log owner]",
            "Date range: [start to end]",
            "Key finding: [largest contributor or top-three share]",
            "Next action: [Fishbone, 5 Whys, or containment check]",
            "Prepared by/date: [name and date]",
        ),
    ),
    "check_sheet": TemplateContent(
        title="Check Sheet",
        stage_label="Understand Current Condition",
        purpose="Purpose: collect defects or events consistently before analysis.",
        data_title="DATA ENTRY - replace check items and tallies only",
        data_columns=("Check item", "Mon", "Tue", "Wed", "Thu", "Fri", "Total"),
        data_rows=(
            ("Wrong label", "8", "7", "9", "10", "8", "42"),
            ("Missing label", "5", "6", "4", "7", "5", "27"),
            ("Damaged barcode", "4", "3", "5", "3", "3", "18"),
            ("Wrong carton", "1", "2", "2", "2", "2", "9"),
        ),
        demo_title="DEMO EXAMPLE - not project evidence",
        demo_points=(
            "Change data only: replace the check items and daily tallies.",
            "Check item: {{check_item}}",
            "Tally area: {{tally_area}}",
            "Collection notes: {{review_notes}}",
        ),
        project_title="Blank working slide or worksheet",
        project_points=(
            "Project title: [process, location, and period]",
            "Worksheet: replace demo rows with project check items.",
            "Source: [collector, area, system, or observation point]",
            "Date range: [start to end]",
            "Key conclusion: [what the tally makes visible]",
            "Next action: [Pareto, stratification, or follow-up check]",
            "Prepared by/date: [name and date]",
        ),
    ),
    "5w2h": TemplateContent(
        title="5W2H",
        stage_label="Define Problem",
        purpose="Purpose: turn a broad issue into a concrete problem statement.",
        data_title="DATA ENTRY - replace 5W2H answers only",
        data_columns=("Field", "Project answer"),
        data_rows=(
            ("What", "Packing labels are applied incorrectly"),
            ("Why", "Creates rework and shipment delay"),
            ("Where", "Final packing station"),
            ("When", "Baseline week 2026-07-01"),
            ("Who", "Packing team and QC reviewer"),
            ("How", "Wrong SKU label selected before boxing"),
            ("How much", "42 wrong-label cases in baseline week"),
        ),
        demo_title="DEMO EXAMPLE - not project evidence",
        demo_points=(
            "Change data only: replace the answer cells.",
            "What: {{what}}",
            "Why: {{why}}",
            "Where: {{where}}",
            "When: {{when}}",
            "Who: {{who}}",
            "How: {{how}}",
            "How much: {{how_much}}",
        ),
        project_title="Blank working slide or worksheet",
        project_points=(
            "Project title: [problem statement draft]",
            "Worksheet: replace demo answer cells with project facts.",
            "Source: [observation, record, interview, or gemba walk]",
            "Date range: [when the facts were observed]",
            "Key conclusion: [specific measurable problem statement]",
            "Next action: [Check Sheet, Pareto, SIPOC, or target setting]",
            "Prepared by/date: [name and date]",
        ),
    ),
    "fishbone_diagram": TemplateContent(
        title="Fishbone Diagram",
        stage_label="Analyze Causes",
        purpose="Purpose: organize suspected causes before verification.",
        data_title="DATA ENTRY - replace cause statements only",
        data_columns=("Branch", "Suspected cause", "Evidence / verification"),
        data_rows=(
            ("Method", "Label check happens after boxing", "Observed in flow walk"),
            ("Machine", "Scanner warning is easy to miss", "Needs confirmation"),
            ("Material", "Similar label roll colors", "Photo evidence needed"),
            ("People", "New operators rotate into station", "Training record review"),
            ("Environment", "Work table has mixed SKU paperwork", "Gemba note"),
        ),
        demo_title="DEMO EXAMPLE - not project evidence",
        demo_points=(
            "Change data only: replace branch causes and evidence notes.",
            "Effect statement: {{effect_statement}}",
            "Cause branches: {{cause_branches}}",
            "Verified causes: {{verified_causes}}",
        ),
        project_title="Blank working slide or worksheet",
        project_points=(
            "Project title: [effect statement]",
            "Diagram: replace demo branches and causes.",
            "Source: [team session, gemba notes, or evidence review]",
            "Date range: [problem period or session date]",
            "Key conclusion: [which causes need verification]",
            "Next action: [5 Whys, data check, or verification plan]",
            "Prepared by/date: [name and date]",
        ),
    ),
    "5_whys": TemplateContent(
        title="5 Whys",
        stage_label="Analyze Causes",
        purpose="Purpose: trace a problem through evidence-backed why links.",
        data_title="DATA ENTRY - replace why-chain answers only",
        data_columns=("Step", "Why answer", "Evidence check"),
        data_rows=(
            ("Problem", "Wrong label applied to packed cartons", "Pareto category"),
            ("Why 1", "Operator selected wrong label roll", "Observed case review"),
            ("Why 2", "Similar rolls stored together", "Storage check"),
            ("Why 3", "No visual separation standard", "SOP review"),
            ("Why 4", "Storage standard was never updated", "Owner interview"),
            ("Why 5", "Change-control step missing", "Procedure audit"),
        ),
        demo_title="DEMO EXAMPLE - not project evidence",
        demo_points=(
            "Change data only: replace the why-chain answers and evidence checks.",
            "Problem: {{problem}}",
            "Why chain: {{why_chain}}",
            "Root cause: {{root_cause}}",
            "Verification: {{verification}}",
        ),
        project_title="Blank working slide or worksheet",
        project_points=(
            "Project title: [selected problem or cause]",
            "Worksheet: replace demo why-chain answers.",
            "Source: [observation, interview, record, or verification check]",
            "Date range: [when evidence was checked]",
            "Key conclusion: [candidate root cause and evidence status]",
            "Next action: [verify cause or design countermeasure]",
            "Prepared by/date: [name and date]",
        ),
    ),
}


TEMPLATE_NATIVE_GUIDANCE: dict[str, MethodKitGuidance] = {
    "check_sheet": MethodKitGuidance(
        qcc_stage_fit=(
            "Understand Current Condition: collect baseline facts.",
            "Use before Pareto, Histogram, or stratification.",
            "A tally shows occurrence, not root cause.",
        ),
        required_inputs=(
            "Clear observation definition.",
            "Check items or categories.",
            "Collection location and collector.",
            "Source, date range, filters, and notes.",
        ),
        when_to_use=(
            "The team needs consistent manual counts.",
            "Categories, shifts, locations, or dates must be compared.",
            "Data collection rules can be explained simply.",
        ),
        when_not_to_use=(
            "Observers disagree about what counts.",
            "Categories overlap or change midstream.",
            "The team needs final calculation proof without source data.",
        ),
        interpretation_patterns=(
            "Most frequent check item: name the item and count.",
            "Pattern: compare location, shift, or date differences.",
            "Next action: convert validated tallies into Pareto or trend data.",
            "Caution: the sheet does not prove root cause.",
        ),
        common_mistakes=(
            "Changing categories during collection.",
            "Missing source or date range.",
            "Counting different event definitions together.",
            "Treating demo data as project evidence.",
        ),
        facilitator_checklist=(
            "Observation definition is clear.",
            "Source and date range shown.",
            "Categories are understandable and stable.",
            "Key conclusion and next action are written.",
        ),
        python_assist_decision=(
            "PowerPoint is enough for teaching, draft, and small tally sheets.",
            "Python is not normally needed for this worksheet.",
            "Use Python later only to summarize large exported logs.",
            "High-rigor data claims need a validated path.",
        ),
        evidence_levels=(
            "Level 1: teaching/draft PowerPoint edits.",
            "Level 2: normal project with source and checklist evidence.",
            "Level 3: formal review with source data and versioned template.",
            "Level 4: audit/high-risk reproducible evidence.",
        ),
        evidence_source_note=(
            "Record source.",
            "Record date range.",
            "Record collection rules and filters.",
            "Record collector and review notes.",
        ),
    ),
    "5w2h": MethodKitGuidance(
        qcc_stage_fit=(
            "Define Problem: turn broad concern into project scope.",
            "Use before detailed data collection and target setting.",
            "The worksheet frames facts, not root causes.",
        ),
        required_inputs=(
            "Observed problem background.",
            "Where, when, and who context.",
            "How the issue is detected or handled.",
            "Source, date range, and measurable impact.",
        ),
        when_to_use=(
            "The problem statement is vague.",
            "The team needs shared scope before analysis.",
            "Known facts can be separated from assumptions.",
        ),
        when_not_to_use=(
            "The team is already choosing countermeasures.",
            "Facts are unknown and not marked as gaps.",
            "The slide is being used as root-cause proof.",
        ),
        interpretation_patterns=(
            "Problem statement: combine What, Where, When, and How much.",
            "Scope decision: state what is included and excluded.",
            "Gap: identify unknown facts for Check Sheet collection.",
            "Caution: do not include unverified causes.",
        ),
        common_mistakes=(
            "Writing opinions instead of observations.",
            "Embedding the solution in the problem statement.",
            "Missing measurable impact.",
            "Treating demo content as project evidence.",
        ),
        facilitator_checklist=(
            "Each W/H answer is specific.",
            "Source and date range shown where facts are cited.",
            "No countermeasure is hidden in the problem statement.",
            "Key conclusion and next action are written.",
        ),
        python_assist_decision=(
            "PowerPoint is enough for teaching, draft, and normal 5W2H use.",
            "Python is not normally needed for this reasoning worksheet.",
            "Use Python later only for evidence summaries behind impact numbers.",
            "High-rigor data claims need a validated path.",
        ),
        evidence_levels=(
            "Level 1: teaching/draft PowerPoint edits.",
            "Level 2: normal project with source and checklist evidence.",
            "Level 3: formal review with source data and versioned template.",
            "Level 4: audit/high-risk reproducible evidence.",
        ),
        evidence_source_note=(
            "Record source.",
            "Record date range.",
            "Record known facts and assumptions.",
            "Record prepared by/date.",
        ),
    ),
    "fishbone_diagram": MethodKitGuidance(
        qcc_stage_fit=(
            "Analyze Causes: organize suspected causes for one effect.",
            "Use after the problem and baseline evidence are clear.",
            "The diagram is a hypothesis map, not proof.",
        ),
        required_inputs=(
            "One clear effect statement.",
            "Cause branches that fit the process.",
            "Suspected causes and evidence notes.",
            "Source, date range, and verification plan.",
        ),
        when_to_use=(
            "The team needs to organize many possible causes.",
            "Several process factors may contribute to the effect.",
            "Follow-up verification can be assigned.",
        ),
        when_not_to_use=(
            "The effect statement is vague or mixed.",
            "The team wants to vote without evidence.",
            "The slide is being used as verified root-cause proof.",
        ),
        interpretation_patterns=(
            "Cause map: name the highest-priority branches.",
            "Verification decision: list causes to test next.",
            "Next action: run 5 Whys or data checks for selected causes.",
            "Caution: suspected causes remain unverified until tested.",
        ),
        common_mistakes=(
            "Mixing several effects on one diagram.",
            "Listing symptoms as causes.",
            "Marking causes verified without evidence.",
            "Treating demo content as project evidence.",
        ),
        facilitator_checklist=(
            "Effect statement matches project evidence.",
            "Source and date range or session date shown.",
            "Suspected and verified causes are separated.",
            "Key conclusion and next action are written.",
        ),
        python_assist_decision=(
            "PowerPoint is enough for teaching, draft, and diagram editing.",
            "Python is not normally needed for this visual method.",
            "Use Python later only to summarize verification data.",
            "High-rigor data claims need a validated path.",
        ),
        evidence_levels=(
            "Level 1: teaching/draft PowerPoint edits.",
            "Level 2: normal project with source and checklist evidence.",
            "Level 3: formal review with source data and versioned template.",
            "Level 4: audit/high-risk reproducible evidence.",
        ),
        evidence_source_note=(
            "Record source.",
            "Record date range or session date.",
            "Record assumptions and verification status.",
            "Record prepared by/date.",
        ),
    ),
    "5_whys": MethodKitGuidance(
        qcc_stage_fit=(
            "Analyze Causes: deepen one selected problem or cause.",
            "Use after 5W2H, Pareto, or Fishbone narrows focus.",
            "A why chain needs evidence before countermeasures.",
        ),
        required_inputs=(
            "Narrow starting problem.",
            "Why answers and evidence checks.",
            "Candidate root cause.",
            "Source, date range, and verification status.",
        ),
        when_to_use=(
            "The team can follow one cause chain.",
            "Each answer can be checked.",
            "The output will guide verification or countermeasures.",
        ),
        when_not_to_use=(
            "The problem has many independent branches.",
            "Answers are opinions with no evidence path.",
            "The fifth answer is assumed true without verification.",
        ),
        interpretation_patterns=(
            "Cause chain: summarize the evidence-backed path.",
            "Root-cause candidate: name the cause and verification status.",
            "Next action: verify the link or design a countermeasure.",
            "Caution: the worksheet does not prove root cause by count of whys.",
        ),
        common_mistakes=(
            "Branching into multiple problems.",
            "Skipping evidence checks.",
            "Stopping at a convenient answer.",
            "Treating demo content as project evidence.",
        ),
        facilitator_checklist=(
            "Starting problem is narrow.",
            "Source and date range shown for evidence checks.",
            "Each why follows from the previous answer.",
            "Key conclusion and next action are written.",
        ),
        python_assist_decision=(
            "PowerPoint is enough for teaching, draft, and normal 5 Whys use.",
            "Python is not normally needed for this reasoning worksheet.",
            "Use Python later only to summarize verification data.",
            "High-rigor data claims need a validated path.",
        ),
        evidence_levels=(
            "Level 1: teaching/draft PowerPoint edits.",
            "Level 2: normal project with source and checklist evidence.",
            "Level 3: formal review with source data and versioned template.",
            "Level 4: audit/high-risk reproducible evidence.",
        ),
        evidence_source_note=(
            "Record source.",
            "Record date range.",
            "Record evidence checks and assumptions.",
            "Record prepared by/date.",
        ),
    ),
}


def main() -> None:
    catalog = load_template_catalog(CATALOG_PATH)
    for entry in catalog.templates:
        if entry.source_file is None:
            raise ValueError(f"{entry.template_id} is missing source_file.")
        if not Path(entry.source_file).exists():
            raise FileNotFoundError(entry.source_file)
        if not Path(entry.markdown_guide).exists():
            raise FileNotFoundError(entry.markdown_guide)
        _build_template(entry)


def _build_template(entry: TemplateCatalogEntry) -> None:
    content = CONTENT[entry.method_id]
    output_path = Path(entry.file)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = f"{content.title} QCC Method Template"
    prs.core_properties.subject = "QCC Toolkit method template"
    prs.core_properties.keywords = (
        f"qcc, {entry.method_id}, {entry.template_id}, DEMO EXAMPLE"
    )
    prs.core_properties.created = GENERATED_AT
    prs.core_properties.modified = GENERATED_AT
    prs.core_properties.last_modified_by = "QCC Toolkit"
    prs.core_properties.comments = (
        "Generated by tools/build_ppt_templates.py from template catalog and "
        "reviewable Markdown source notes."
    )

    _add_overview_slide(prs, entry, content)
    _add_data_entry_slide(prs, entry, content)
    if entry.method_id == "pareto_chart":
        _add_pareto_method_kit_slides(prs, entry, content)
    else:
        _add_template_native_method_kit_slides(prs, entry, content)
    _add_demo_slide(prs, entry, content)
    _add_project_slide(prs, entry, content)

    prs.save(output_path)
    _normalize_pptx_package(output_path)


def _normalize_pptx_package(path: Path) -> None:
    """Rewrite the zip package with stable entry order and timestamps."""

    with ZipFile(path) as source:
        entries = [(name, source.read(name)) for name in sorted(source.namelist())]

    with ZipFile(path, "w", compression=ZIP_DEFLATED) as target:
        for name, data in entries:
            if name.startswith("ppt/embeddings/") and data.startswith(b"PK"):
                data = _normalize_nested_office_package(data)
            info = ZipInfo(name, ZIP_TIMESTAMP)
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)


def _normalize_nested_office_package(data: bytes) -> bytes:
    """Normalize embedded Office zip packages such as chart workbooks."""

    source_buffer = BytesIO(data)
    target_buffer = BytesIO()
    with ZipFile(source_buffer) as source:
        entries = [(name, source.read(name)) for name in sorted(source.namelist())]

    with ZipFile(target_buffer, "w", compression=ZIP_DEFLATED) as target:
        for name, entry_data in entries:
            if name == "docProps/core.xml":
                entry_data = _normalize_core_properties(entry_data)
            info = ZipInfo(name, ZIP_TIMESTAMP)
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, entry_data)
    return target_buffer.getvalue()


def _normalize_core_properties(data: bytes) -> bytes:
    text = data.decode("utf-8")
    text = _replace_tag_text(text, "dcterms:created", "2026-07-08T00:00:00Z")
    text = _replace_tag_text(text, "dcterms:modified", "2026-07-08T00:00:00Z")
    return text.encode("utf-8")


def _replace_tag_text(text: str, tag: str, value: str) -> str:
    start = text.find(f"<{tag}")
    if start == -1:
        return text
    content_start = text.find(">", start)
    if content_start == -1:
        return text
    end = text.find(f"</{tag}>", content_start)
    if end == -1:
        return text
    return f"{text[: content_start + 1]}{value}{text[end:]}"


def _add_overview_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, content.title, "QCC method template")
    _add_meta_bar(slide, entry, content)
    _add_text_box(
        slide,
        left=0.7,
        top=1.75,
        width=5.6,
        height=1.25,
        text=content.purpose,
        font_size=24,
        bold=True,
        color=RGBColor(35, 45, 55),
    )
    _add_section_box(
        slide,
        left=0.7,
        top=3.35,
        width=5.8,
        height=2.6,
        title="Template rules",
        items=(
            "PowerPoint teaches and presents.",
            "Markdown governs method knowledge.",
            "Python generates data-dependent evidence.",
            "Final project slides must preserve evidence references.",
        ),
    )
    _add_section_box(
        slide,
        left=6.85,
        top=3.35,
        width=5.8,
        height=2.6,
        title="Placeholders",
        items=tuple(f"{{{{{name}}}}}" for name in entry.expected_placeholders),
    )
    _add_footer(slide, entry)


def _add_demo_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, content.demo_title, content.title)
    _add_badge(slide, "DEMO EXAMPLE - not project evidence", 0.7, 1.15, 4.2)
    _add_section_box(
        slide,
        left=0.7,
        top=1.85,
        width=5.9,
        height=4.6,
        title="Demo structure",
        items=content.demo_points,
    )
    _add_placeholder_panel(slide, entry, left=6.95, top=1.85)
    _add_footer(slide, entry)


def _add_data_entry_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, content.data_title, content.title)
    _add_badge(slide, "Change data only - keep layout and labels", 0.7, 1.15, 5.0)
    _add_data_table(
        slide,
        columns=content.data_columns,
        rows=content.data_rows,
        left=0.7,
        top=1.8,
        width=6.0,
        height=4.75,
    )
    if entry.method_id == "pareto_chart":
        _add_pareto_chart(slide, content, left=7.0, top=1.8, width=5.4, height=3.2)
        _add_text_box(
            slide,
            left=7.0,
            top=5.25,
            width=5.4,
            height=0.7,
            text=(
                "Use PowerPoint Edit Data for draft visuals. "
                "Use Python-generated evidence for final QCC conclusions."
            ),
            font_size=10,
            color=RGBColor(80, 88, 96),
            fill=RGBColor(238, 244, 243),
            border=RGBColor(170, 195, 190),
        )
    else:
        _add_section_box(
            slide,
            left=7.0,
            top=1.8,
            width=5.4,
            height=4.75,
            title="How to use this template",
            items=(
                "Replace the sample rows with project data.",
                "Keep the method labels and review fields.",
                "Move the summarized result into the project slide.",
                "Preserve evidence notes for review.",
            ),
        )
    _add_footer(slide, entry)


def _add_pareto_method_kit_slides(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="QCC stage fit and required inputs",
        left_title="QCC stage fit",
        left_items=(
            "Understand Current Condition: rank baseline categories.",
            "Analyze Causes: choose where cause analysis starts.",
            "Pareto shows concentration, not root cause.",
        ),
        right_title="Required inputs",
        right_items=(
            "Category and count.",
            "Non-overlapping categories.",
            "Consistent data period.",
            "Source, date range, filters, and calculation notes.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="When to use and when not to use",
        left_title="When to use",
        left_items=(
            "Defect, complaint, delay, or event categories can be counted.",
            "The team needs a focus decision before deeper analysis.",
            "Small category-count data can be edited in PowerPoint.",
        ),
        right_title="When not to use",
        right_items=(
            "Categories overlap or were changed after seeing the answer.",
            "Data periods are mixed.",
            "The team needs proof of root cause.",
            "The conclusion is high-rigor without reproducible evidence.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="PowerPoint edit instructions",
        left_title="Edit chart data",
        left_items=(
            "Right-click the chart and choose Edit Data.",
            "Replace demo categories and counts.",
            "Sort counts descending.",
            "Verify formula cells were not overwritten.",
        ),
        right_title="Review before use",
        right_items=(
            "Keep source and date range visible.",
            "Do not present demo data as project evidence.",
            "Check cumulative values if used.",
            "Write Key finding and Next action.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Chart decision guide",
        left_title="Decision supported",
        left_items=(
            "Choose the category or vital few for follow-up analysis.",
            "Pattern to look for: dominant category, steep drop, or long tail.",
            "Next method: Fishbone, 5 Whys, stratification, or containment.",
        ),
        right_title="Safe conclusion",
        right_items=(
            "Safe conclusion: counted problems are concentrated here.",
            "Overclaim to avoid: Pareto does not prove root cause.",
            "Do not claim stable improvement from one Pareto snapshot.",
        ),
    )
    _add_pareto_chart_variant_slide(prs, entry, content)
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Chart quality checklist",
        left_title="Chart fields",
        left_items=(
            "Source.",
            "Date range.",
            "Filters.",
            "Percent.",
            "Cumulative percent.",
            "Formula check.",
        ),
        right_title="Reviewer checks",
        right_items=(
            "Bars sorted descending.",
            "Category definitions are clear.",
            "Labels and counts are readable.",
            "Formula cells were not overwritten.",
            "Key finding and next action are written.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Interpretation patterns and common mistakes",
        left_title="Interpretation patterns",
        left_items=(
            "Largest contributor: name the top category and count.",
            "Top-three share: state combined contribution.",
            "Focus decision: name the next method.",
            "Caution: Pareto does not prove root cause.",
        ),
        right_title="Common mistakes",
        right_items=(
            "Unsorted bars.",
            "Overlapping categories.",
            "Mixed time periods.",
            "Missing source or date range.",
            "Demo data used as project evidence.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Facilitator checklist and Python assist decision",
        left_title="Facilitator checklist",
        left_items=(
            "Source and date range shown.",
            "Categories clear and non-overlapping.",
            "Counts sorted descending.",
            "Conclusion and next action written.",
        ),
        right_title="Python assist decision",
        right_items=(
            "PowerPoint: training, draft, and small counted tables.",
            "Python: raw logs, repeated generation, validation-heavy data.",
            "Level 3: Python recommended for raw or repeated chart methods.",
            "Level 4: reproducible evidence package or validated path.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Evidence/source note",
        left_title="Evidence levels",
        left_items=(
            "Level 1: teaching/draft PowerPoint edits.",
            "Level 2: normal project with source and checklist evidence.",
            "Level 3: formal review with source data and calculation table.",
            "Level 4: audit/high-risk reproducible evidence.",
        ),
        right_title="Record on the slide",
        right_items=(
            "Source.",
            "Date range.",
            "Filters and assumptions.",
            "Calculation notes.",
            "Prepared by/date.",
        ),
    )


def _add_pareto_chart_variant_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Chart variant library", content.title)
    _add_pareto_chart(slide, content, left=0.75, top=1.55, width=5.8, height=3.25)
    _add_section_box(
        slide,
        left=6.85,
        top=1.45,
        width=5.7,
        height=2.15,
        title="Pareto variants",
        items=(
            "Simple count Pareto: draft focus decision.",
            "Cumulative Pareto: top-share or threshold wording.",
            "Before/after Pareto comparison: comparable periods only.",
            "Focus annotation: mark selected vital few.",
        ),
    )
    _add_section_box(
        slide,
        left=6.85,
        top=3.85,
        width=5.7,
        height=1.85,
        title="When to use Python assist",
        items=(
            "Raw logs or repeated generation.",
            "Validation-heavy data or many categories.",
            "Final reviewed evidence needing metadata.",
        ),
    )
    _add_section_box(
        slide,
        left=0.75,
        top=5.2,
        width=5.8,
        height=0.9,
        title="Focus annotation",
        items=("Call out the selected vital few and the next method.",),
    )
    _add_footer(slide, entry)


def _add_template_native_method_kit_slides(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    guidance = TEMPLATE_NATIVE_GUIDANCE[entry.method_id]
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="QCC stage fit and required inputs",
        left_title="QCC stage fit",
        left_items=guidance.qcc_stage_fit,
        right_title="Required inputs",
        right_items=guidance.required_inputs,
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="When to use and when not to use",
        left_title="When to use",
        left_items=guidance.when_to_use,
        right_title="When not to use",
        right_items=guidance.when_not_to_use,
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Blank working slide or worksheet",
        left_title="Completed demo example",
        left_items=(
            "Review the completed demo example before editing.",
            "DEMO EXAMPLE - not project evidence.",
            "Copy the blank working slide or worksheet for project use.",
            "Replace demo text with project facts.",
        ),
        right_title="Blank working slide or worksheet",
        right_items=(
            "Keep method labels visible.",
            "Record source and date range.",
            "Write Key conclusion.",
            "Write Next action.",
        ),
    )
    if entry.method_id == "fishbone_diagram":
        _add_fishbone_diagram_quality_slides(prs, entry, content)
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Interpretation patterns and common mistakes",
        left_title="Interpretation patterns",
        left_items=guidance.interpretation_patterns,
        right_title="Common mistakes",
        right_items=guidance.common_mistakes,
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Facilitator checklist and Python assist decision",
        left_title="Facilitator checklist",
        left_items=guidance.facilitator_checklist,
        right_title="Python assist decision",
        right_items=guidance.python_assist_decision,
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Evidence/source note",
        left_title="Evidence levels",
        left_items=guidance.evidence_levels,
        right_title="Evidence/source note",
        right_items=guidance.evidence_source_note,
    )


def _add_fishbone_diagram_quality_slides(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Diagram quality guide",
        left_title="Diagram decision",
        left_items=(
            "Select suspected causes to verify next for one effect.",
            "Pattern to look for: testable causes, not the most opinions.",
            "Good structure: one effect, branches, causes, markers, evidence.",
            "Safe conclusion: suspected causes and verification priorities.",
        ),
        right_title="Overclaim to avoid",
        right_items=(
            "The diagram does not prove root cause.",
            "The diagram does not prove countermeasure effectiveness.",
            "Do not mix several effects on one fishbone.",
            "Do not select causes without a verification method.",
        ),
    )
    _add_two_column_guidance_slide(
        prs,
        entry,
        content,
        title="Verification marker legend and cause wording guide",
        left_title="Verification marker legend",
        left_items=(
            "[S] Suspected: proposed but not selected.",
            "[V?] Selected for verification: check next.",
            "[V] Verified: supported by later evidence.",
            "[X] Rejected: checked and not supported.",
        ),
        right_title="Cause wording guide",
        right_items=(
            'Weak wording: "Operator error".',
            'Testable wording: "New operators lack label-check training".',
            'Weak wording: "Bad machine".',
            'Testable wording: "Scanner warning is not visible".',
        ),
    )
    _add_editable_fishbone_slide(prs, entry, content)
    _add_cause_verification_plan_slide(prs, entry, content)


def _add_editable_fishbone_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Editable fishbone diagram", content.title)
    _add_badge(slide, "Replace demo causes with project evidence", 0.7, 1.15, 5.4)

    _add_text_box(
        slide,
        left=9.45,
        top=2.7,
        width=2.75,
        height=0.75,
        text="Effect statement\n{{effect_statement}}",
        font_size=13,
        bold=True,
        color=RGBColor(27, 39, 51),
        fill=RGBColor(238, 244, 243),
        border=RGBColor(66, 126, 120),
    )
    _add_line(slide, 1.05, 3.08, 9.45, 3.08, RGBColor(66, 126, 120), 2.25)

    branches = (
        ("Method", "S: Label check after boxing", 1.25, 2.75, 1.35, True),
        ("Machine", "V?: Scanner warning easy to miss", 3.55, 4.9, 1.35, True),
        ("Material", "S: Similar label roll colors", 5.75, 7.1, 1.35, True),
        ("People", "S: New operators rotate in", 1.35, 2.85, 4.55, False),
        ("Environment", "V?: Mixed SKU paperwork", 3.75, 5.25, 4.55, False),
        ("Measurement", "X: Audit form not root cause", 6.05, 7.55, 4.55, False),
    )
    for branch, cause, x1, x2, y2, upper in branches:
        _add_line(slide, x1, 3.08, x2, y2, RGBColor(120, 137, 134), 1.5)
        text_top = y2 - 0.68 if upper else y2 + 0.1
        _add_text_box(
            slide,
            left=x2 - 0.62,
            top=text_top,
            width=2.05,
            height=0.72,
            text=f"{branch}\n{cause}",
            font_size=9,
            bold=False,
            color=RGBColor(35, 45, 55),
            fill=RGBColor(248, 249, 247),
            border=RGBColor(190, 196, 191),
        )

    _add_section_box(
        slide,
        left=0.7,
        top=5.85,
        width=4.0,
        height=0.85,
        title="Verification marker legend",
        items=(
            "[S] Suspected | [V?] Selected for verification | "
            "[V] Verified | [X] Rejected",
        ),
    )
    _add_section_box(
        slide,
        left=4.95,
        top=5.85,
        width=3.55,
        height=0.85,
        title="Selected causes to verify",
        items=("Choose 2-3 causes with clear verification methods.",),
    )
    _add_section_box(
        slide,
        left=8.75,
        top=4.55,
        width=3.65,
        height=1.65,
        title="Evidence/source fields",
        items=(
            "Source: [team session, gemba, records].",
            "Date range/session date: [date].",
            "Prepared by/date: [name/date].",
        ),
    )
    _add_footer(slide, entry)


def _add_cause_verification_plan_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Cause verification plan", content.title)
    _add_data_table(
        slide,
        columns=(
            "Selected cause",
            "Verification method",
            "Owner",
            "Due date",
            "Status",
        ),
        rows=(
            (
                "Scanner warning easy to miss",
                "Gemba observation and photo check",
                "[owner]",
                "[date]",
                "[V?]",
            ),
            (
                "Mixed SKU paperwork",
                "Review packing station layout",
                "[owner]",
                "[date]",
                "[V?]",
            ),
            (
                "[cause]",
                "[check, data review, or 5 Whys]",
                "[owner]",
                "[date]",
                "[S/V?/V/X]",
            ),
        ),
        left=0.75,
        top=1.55,
        width=11.85,
        height=2.2,
    )
    _add_section_box(
        slide,
        left=0.75,
        top=4.05,
        width=5.65,
        height=1.55,
        title="Use this plan",
        items=(
            "Select only causes that can be checked.",
            "Assign owner and due date before countermeasure selection.",
            "Update Status after verification evidence is reviewed.",
        ),
    )
    _add_section_box(
        slide,
        left=6.85,
        top=4.05,
        width=5.65,
        height=1.55,
        title="Review risks",
        items=(
            "No verification method.",
            "Owner or due date missing.",
            "Status marked [V] without evidence.",
            "Cause moved to countermeasure too early.",
        ),
    )
    _add_footer(slide, entry)


def _add_two_column_guidance_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
    *,
    title: str,
    left_title: str,
    left_items: tuple[str, ...],
    right_title: str,
    right_items: tuple[str, ...],
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, title, content.title)
    _add_section_box(
        slide,
        left=0.7,
        top=1.45,
        width=5.7,
        height=4.95,
        title=left_title,
        items=left_items,
    )
    _add_section_box(
        slide,
        left=6.85,
        top=1.45,
        width=5.7,
        height=4.95,
        title=right_title,
        items=right_items,
    )
    _add_footer(slide, entry)


def _add_project_slide(
    prs: Presentation,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, content.project_title, content.title)
    _add_badge(slide, "Replace demo content with project evidence", 0.7, 1.15, 5.2)
    _add_section_box(
        slide,
        left=0.7,
        top=1.85,
        width=5.9,
        height=4.6,
        title="Project slide checklist",
        items=content.project_points,
    )
    _add_placeholder_panel(slide, entry, left=6.95, top=1.85)
    _add_footer(slide, entry)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_header(slide, title: str, subtitle: str) -> None:
    _add_text_box(
        slide,
        left=0.55,
        top=0.28,
        width=9.1,
        height=0.46,
        text=title,
        font_size=24,
        bold=True,
        color=RGBColor(27, 39, 51),
    )
    _add_text_box(
        slide,
        left=0.57,
        top=0.78,
        width=8.5,
        height=0.26,
        text=subtitle,
        font_size=11,
        color=RGBColor(90, 99, 108),
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(1.03),
        Inches(12.2),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(66, 126, 120)
    line.line.fill.background()


def _add_meta_bar(
    slide,
    entry: TemplateCatalogEntry,
    content: TemplateContent,
) -> None:
    text = (
        f"template_id: {entry.template_id}\n"
        f"method_id: {entry.method_id}\n"
        f"qcc_stage: {content.stage_label}\n"
        f"{{{{method_name}}}} | {{{{qcc_stage}}}} | {{{{demo_label}}}}"
    )
    _add_text_box(
        slide,
        left=7.05,
        top=1.45,
        width=5.15,
        height=1.25,
        text=text,
        font_size=12,
        color=RGBColor(45, 55, 65),
        fill=RGBColor(238, 244, 243),
        border=RGBColor(170, 195, 190),
    )


def _add_section_box(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    items: tuple[str, ...],
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 249, 247)
    shape.line.color.rgb = RGBColor(190, 196, 191)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.18)
    frame.margin_top = Inches(0.14)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(27, 39, 51)
    for item in items:
        p = frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(45, 55, 65)


def _add_data_table(
    slide,
    *,
    columns: tuple[str, ...],
    rows: tuple[tuple[str, ...], ...],
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    for column_idx, column in enumerate(columns):
        cell = table.cell(0, column_idx)
        cell.text = column
        _format_cell(cell, fill=RGBColor(66, 126, 120), color=RGBColor(255, 255, 255))

    for row_idx, row in enumerate(rows, start=1):
        for column_idx, value in enumerate(row):
            cell = table.cell(row_idx, column_idx)
            cell.text = value
            fill = RGBColor(248, 249, 247) if row_idx % 2 else RGBColor(238, 244, 243)
            _format_cell(cell, fill=fill, color=RGBColor(35, 45, 55))


def _format_cell(cell, *, fill: RGBColor, color: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(8)
        paragraph.font.color.rgb = color


def _add_pareto_chart(
    slide,
    content: TemplateContent,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = tuple(row[0] for row in content.data_rows)
    chart_data.add_series("Count", tuple(int(row[1]) for row in content.data_rows))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)


def _add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor,
    width: float,
) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)


def _add_placeholder_panel(
    slide,
    entry: TemplateCatalogEntry,
    *,
    left: float,
    top: float,
) -> None:
    _add_section_box(
        slide,
        left=left,
        top=top,
        width=5.35,
        height=4.6,
        title="Reserved placeholders",
        items=tuple(f"{{{{{name}}}}}" for name in entry.expected_placeholders),
    )


def _add_badge(slide, text: str, left: float, top: float, width: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.36),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 242, 204)
    shape.line.color.rgb = RGBColor(214, 159, 46)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(102, 75, 17)


def _add_footer(slide, entry: TemplateCatalogEntry) -> None:
    _add_text_box(
        slide,
        left=0.6,
        top=6.95,
        width=11.9,
        height=0.25,
        text=(
            f"{entry.template_id} | {entry.method_id} | "
            "Final data-dependent evidence belongs in the evidence package."
        ),
        font_size=8,
        color=RGBColor(105, 112, 119),
    )


def _add_text_box(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    fill: RGBColor | None = None,
    border: RGBColor | None = None,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if border is not None:
        box.line.color.rgb = border
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


if __name__ == "__main__":
    main()
